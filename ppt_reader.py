from pptx import Presentation


def extract_text(file):
    presentation = Presentation(file)

    slides = []

    for slide in presentation.slides:
        slide_text = []

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()

                if text:
                    slide_text.append(text)

        slides.append("\n".join(slide_text))

    return slides